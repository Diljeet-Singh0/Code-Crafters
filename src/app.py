from flask import Flask, render_template, request, send_file
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches
import os

app = Flask(__name__)

# Route for the index page
@app.route('/')
def index():
    return render_template('index.html')

# Route for the form
@app.route('/form', methods=['POST'])
def form():
    doc_type = request.form['doc_type']
    return render_template('form.html', doc_type=doc_type)

# Route to generate report
@app.route('/generate_report', methods=['POST'])
def generate_report():
    data = {
        "name": request.form.get('name', ''),
        "class": request.form.get('class', ''),
        "college_roll_no": request.form.get('college_roll_no', ''),
        "university_roll_no": request.form.get('university_roll_no', ''),
        "contact_number": request.form.get('contact_number', ''),
        "fathers_name": request.form.get('fathers_name', ''),
        "mothers_name": request.form.get('mothers_name', ''),
        "address": request.form.get('address', '')
    }
    
    report_path = 'report.pdf'
    
    # Create a PDF report
    c = canvas.Canvas(report_path, pagesize=letter)
    c.drawString(100, 750, "Guru Nanak Dev Engineering College")
    y_position = 730
    for key, value in data.items():
        if value:  # Only display filled fields
            c.drawString(100, y_position, f"{key.replace('_', ' ').title()}: {value}")
            y_position -= 20
    c.save()
    
    return send_file(report_path, as_attachment=True)

# Route to generate presentation
@app.route('/generate_presentation', methods=['POST'])
def generate_presentation():
    data = {
        "name": request.form.get('name', ''),
        "class": request.form.get('class', ''),
        "college_roll_no": request.form.get('college_roll_no', ''),
        "university_roll_no": request.form.get('university_roll_no', ''),
        "contact_number": request.form.get('contact_number', ''),
        "fathers_name": request.form.get('fathers_name', ''),
        "mothers_name": request.form.get('mothers_name', ''),
        "address": request.form.get('address', '')
    }
    
    # Create a PowerPoint presentation
    ppt_path = 'presentation.pptx'
    presentation = Presentation()
    
    # Add a title slide
    slide_layout = presentation.slide_layouts[0]  # 0 is the layout for title slide
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Guru Nanak Dev Engineering College"
    subtitle.text = "Student Information"

    # Add a slide for each piece of data
    for key, value in data.items():
        if value:  # Only add slides for filled fields
            slide_layout = presentation.slide_layouts[1]  # 1 is the layout for content slide
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = key.replace('_', ' ').title()
            content.text = value

    # Save the presentation
    presentation.save(ppt_path)
    
    return send_file(ppt_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)